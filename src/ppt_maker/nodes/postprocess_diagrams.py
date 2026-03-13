from __future__ import annotations

import os
import random
import re
from typing import Any, Dict, List, Optional, Set, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from .gemini_diagram_images import maybe_insert_generated_diagrams

BG_MARKER_NAME = "__RandiBgImage__"
BG_MARKER_ALT = "RandiBackground"

FORMAL_ENDING_RE = re.compile(r"(?:입니다|합니다|있습니다|했습니다)(?=\s*(?:$|[.!?|,;]))")


# 발표 순서 고정
SECTION_ORDER = [
    "기관 소개",
    "연구 개요",
    "연구 필요성",
    "연구 목표",
    "연구 내용",
    "추진 계획",
    "활용방안 및 기대효과",
    "사업화 전략 및 계획",
]


def _norm(s: Any) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()


def _remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _send_shape_to_back(shape) -> None:
    sp = shape._element
    parent = sp.getparent()
    parent.remove(sp)
    # Keep spTree required first nodes intact and insert just after them.
    parent.insert(2, sp)


def _apply_background_image_to_all_slides(prs: Presentation, image_path: str) -> int:
    bg = str(image_path or "").strip()
    if not bg or not os.path.exists(bg):
        return 0

    inserted = 0
    for slide in prs.slides:
        for sh in list(slide.shapes):
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and (
                str(getattr(sh, "name", "")) == BG_MARKER_NAME
                or str(getattr(sh, "alternative_text", "")) == BG_MARKER_ALT
            ):
                _remove_shape(sh)
        pic = slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
        try:
            pic.name = BG_MARKER_NAME
        except Exception:
            pass
        try:
            pic.alternative_text = BG_MARKER_ALT
        except Exception:
            pass
        _send_shape_to_back(pic)
        inserted += 1
    return inserted


def _section_at(deck_json: Dict[str, Any], idx: int) -> str:
    slides = (deck_json or {}).get("slides") or []
    if 0 <= idx < len(slides) and isinstance(slides[idx], dict):
        return _norm(slides[idx].get("section"))
    return ""


def _prompt_type_at(deck_json: Dict[str, Any], idx: int) -> str:
    slides = (deck_json or {}).get("slides") or []
    if 0 <= idx < len(slides) and isinstance(slides[idx], dict):
        return _norm(slides[idx].get("image_prompt_type")).lower()
    return ""


def _resolve_profile_assets(profile: str, state: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    st = state or {}
    module_dir = os.path.dirname(os.path.abspath(__file__))
    default_bg_dir = os.path.normpath(os.path.join(module_dir, "..", "background"))
    base_dir = st.get("postprocess_background_base_dir") or default_bg_dir
    profile_n = _norm(profile).lower()
    assets: Dict[str, Any] = {"profile": profile_n, "base_dir": base_dir}

    if profile_n == "basic":
        assets["main"] = str(
            st.get("postprocess_bg_basic_main")
            or os.path.join(base_dir, "basic_main.png")
        )
        assets["image"] = str(
            st.get("postprocess_bg_basic_image")
            or os.path.join(base_dir, "basic_image.png")
        )
        assets["origin"] = str(
            st.get("postprocess_bg_basic_origin")
            or os.path.join(base_dir, "basic_origin.png")
        )
        return assets

    if profile_n == "brown":
        assets["content"] = str(
            st.get("postprocess_bg_brown_content")
            or os.path.join(base_dir, "brown", "b_content.png")
        )
        assets["origins"] = [
            str(st.get("postprocess_bg_brown_origin1") or os.path.join(base_dir, "brown", "b_origin1.png")),
            str(st.get("postprocess_bg_brown_origin2") or os.path.join(base_dir, "brown", "b_origin2.png")),
            str(st.get("postprocess_bg_brown_origin3") or os.path.join(base_dir, "brown", "b_origin3.png")),
        ]
        return assets

    return assets


def _pick_profile_background_path(
    idx: int,
    total: int,
    section: str,
    prompt_type: str,
    assets: Dict[str, Any],
    *,
    rng: random.Random,
) -> str:
    profile = assets.get("profile")
    sec = _norm(section)

    if profile == "basic":
        main = str(assets.get("main") or "")
        image = str(assets.get("image") or "")
        origin = str(assets.get("origin") or "")
        if idx in {0, total - 1}:
            return main
        if prompt_type in {"system_architecture", "plan_orgchart_fixed"}:
            return image or origin
        return origin

    if profile == "brown":
        content = str(assets.get("content") or "")
        origins = [p for p in (assets.get("origins") or []) if str(p).strip()]
        if idx == 1 or sec == "목차":
            return content
        if origins:
            return rng.choice(origins)
        return ""

    return ""


def _apply_background_profile_to_all_slides(
    prs: Presentation,
    deck_json: Dict[str, Any],
    *,
    profile: str,
    state: Optional[Dict[str, Any]] = None,
) -> int:
    assets = _resolve_profile_assets(profile, state)
    seed_raw = (
        (state or {}).get("postprocess_background_random_seed")
        or os.environ.get("POSTPROCESS_BACKGROUND_RANDOM_SEED")
        or "42"
    )
    try:
        seed = int(seed_raw)
    except Exception:
        seed = 42
    rng = random.Random(seed)

    inserted = 0
    per_file: Dict[str, int] = {}
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides):
        path = _pick_profile_background_path(
            idx,
            total,
            _section_at(deck_json, idx),
            _prompt_type_at(deck_json, idx),
            assets,
            rng=rng,
        )
        if not path or not os.path.exists(path):
            continue
        pic = slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        try:
            pic.name = BG_MARKER_NAME
        except Exception:
            pass
        try:
            pic.alternative_text = BG_MARKER_ALT
        except Exception:
            pass
        _send_shape_to_back(pic)
        inserted += 1
        per_file[path] = int(per_file.get(path, 0)) + 1

    if inserted:
        print(f"[INFO] background profile applied: profile={profile}, seed={seed}, inserted={inserted}")
        for p, n in per_file.items():
            print(f"[INFO] background file usage: {os.path.basename(p)} x {n}")
    else:
        print(f"[WARN] background profile applied but no files inserted: profile={profile}, assets={assets}")
    return inserted


def _remove_background_images(prs: Presentation) -> int:
    removed = 0
    for slide in prs.slides:
        for sh in list(slide.shapes):
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if str(getattr(sh, "name", "")) == BG_MARKER_NAME or str(getattr(sh, "alternative_text", "")) == BG_MARKER_ALT:
                _remove_shape(sh)
                removed += 1
    return removed


def _delete_slide(prs: Presentation, slide_index: int) -> None:
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    sldId = sldIdLst[slide_index]
    rId = sldId.rId
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def _clear_slide(slide) -> None:
    for sh in list(slide.shapes):
        _remove_shape(sh)


def _extract_best_title(state: Dict[str, Any]) -> str:
    default_title = "국가 R&D 제안 발표자료"
    deck = (state.get("deck_json") or {})
    title = _norm(deck.get("deck_title") or "")
    if title and title != "(과제명 미기재)":
        return title
    return default_title


def _write_cover(slide, title: str) -> None:
    _clear_slide(slide)

    # Cover accent lines (requested style)
    _add_solid_rect(slide, left=0.0, top=0.0, width=13.333, height=0.12, rgb=(15, 76, 129))
    _add_solid_rect(slide, left=0.0, top=7.38, width=13.333, height=0.12, rgb=(42, 157, 143))
    _add_solid_rect(slide, left=0.9, top=1.95, width=4.2, height=0.08, rgb=(191, 217, 238))

    # Move title slightly down for better balance.
    tx = slide.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.0))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(42)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)

    # Cover subtitle/tag removed by request.


def _write_agenda(slide) -> None:
    _clear_slide(slide)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_t = title_box.text_frame
    tf_t.clear()
    p0 = tf_t.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "목차"
    r0.font.size = Pt(32)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(20, 20, 20)

    body = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(5.2))
    tf = body.text_frame
    tf.clear()

    for i, item in enumerate(SECTION_ORDER, 1):
        p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
        p.text = f"{i}. {item}"
        p.level = 0
        if p.runs:
            p.runs[0].font.size = Pt(22)
        else:
            r = p.add_run()
            r.font.size = Pt(22)

def _add_solid_rect(slide, *, left: float, top: float, width: float, height: float, rgb: Tuple[int, int, int]) -> None:
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*rgb)
    shp.line.fill.background()


def _decorate_cover_slide(slide) -> None:
    # 상단/하단 컬러 밴드로 표지 배경을 보강한다.
    _add_solid_rect(slide, left=0.0, top=0.0, width=13.333, height=0.45, rgb=(15, 76, 129))
    _add_solid_rect(slide, left=0.0, top=7.2, width=13.333, height=0.18, rgb=(42, 157, 143))
    _add_solid_rect(slide, left=0.8, top=1.4, width=4.6, height=0.16, rgb=(42, 157, 143))
    _add_solid_rect(slide, left=8.8, top=0.62, width=3.7, height=0.10, rgb=(191, 217, 238))
    _add_solid_rect(slide, left=9.6, top=0.86, width=2.9, height=0.10, rgb=(221, 234, 247))
    _add_solid_rect(slide, left=0.9, top=5.1, width=3.1, height=0.10, rgb=(191, 217, 238))


def _decorate_thanks_slide(slide) -> None:
    # 마지막 슬라이드에 동일한 색상 포인트를 추가한다.
    _add_solid_rect(slide, left=0.0, top=0.0, width=0.22, height=7.5, rgb=(15, 76, 129))
    _add_solid_rect(slide, left=0.22, top=7.05, width=13.11, height=0.30, rgb=(232, 241, 250))


def _decorate_content_slides(prs: Presentation) -> None:
    # Apply subtle, consistent template accents on all body slides.
    for idx, slide in enumerate(prs.slides):
        if idx in {0, len(prs.slides) - 1}:  # cover / thanks
            continue
        _add_solid_rect(slide, left=0.0, top=0.0, width=13.333, height=0.08, rgb=(225, 235, 246))
        _add_solid_rect(slide, left=0.0, top=7.42, width=13.333, height=0.08, rgb=(225, 235, 246))
        _add_solid_rect(slide, left=0.0, top=0.0, width=0.08, height=7.5, rgb=(191, 217, 238))


def _style_tables(prs: Presentation) -> None:
    # 헤더/본문 색상을 통일해 가독성을 보강한다.
    for slide in prs.slides:
        for sh in slide.shapes:
            if not getattr(sh, "has_table", False):
                continue
            tbl = sh.table
            rows = len(tbl.rows)
            cols = len(tbl.columns)
            if rows <= 0 or cols <= 0:
                continue

            for r in range(rows):
                for c in range(cols):
                    cell = tbl.cell(r, c)
                    cell.fill.solid()
                    if r == 0:
                        cell.fill.fore_color.rgb = RGBColor(34, 74, 122)    # header
                    elif r % 2 == 1:
                        cell.fill.fore_color.rgb = RGBColor(241, 246, 252)  # zebra1
                    else:
                        cell.fill.fore_color.rgb = RGBColor(250, 252, 255)  # zebra2

                    tf = cell.text_frame
                    tf.word_wrap = True
                    for p in tf.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else RGBColor(28, 33, 39)
                            if r == 0:
                                run.font.bold = True
                                run.font.size = Pt(11)
                            else:
                                # 본문은 폰트를 조금 줄여 overflow를 완화한다.
                                run.font.size = Pt(10)


def _strip_formal_endings_text(text: str) -> str:
    if not text:
        return ""
    out_lines: List[str] = []
    for line in str(text).splitlines():
        t = FORMAL_ENDING_RE.sub("", line)
        t = re.sub(r"\s{2,}", " ", t).strip()
        out_lines.append(t)
    return "\n".join(out_lines)


def _strip_formal_endings_in_presentation(prs: Presentation) -> None:
    for slide in prs.slides:
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                tf = sh.text_frame
                for p in tf.paragraphs:
                    if p.runs:
                        for run in p.runs:
                            run.text = _strip_formal_endings_text(run.text or "")
                    else:
                        p.text = _strip_formal_endings_text(p.text or "")
            if getattr(sh, "has_table", False):
                tbl = sh.table
                for r in range(len(tbl.rows)):
                    for c in range(len(tbl.columns)):
                        cell = tbl.cell(r, c)
                        tf = cell.text_frame
                        for p in tf.paragraphs:
                            if p.runs:
                                for run in p.runs:
                                    run.text = _strip_formal_endings_text(run.text or "")
                            else:
                                p.text = _strip_formal_endings_text(p.text or "")


def _apply_font_name(prs: Presentation, font_name: str) -> None:
    fn = _norm(font_name)
    if not fn:
        return
    for slide in prs.slides:
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                tf = sh.text_frame
                for p in tf.paragraphs:
                    if p.runs:
                        for run in p.runs:
                            run.font.name = fn
                    else:
                        r = p.add_run()
                        r.font.name = fn
            if getattr(sh, "has_table", False):
                tbl = sh.table
                for r_idx in range(len(tbl.rows)):
                    for c_idx in range(len(tbl.columns)):
                        cell = tbl.cell(r_idx, c_idx)
                        for p in cell.text_frame.paragraphs:
                            if p.runs:
                                for run in p.runs:
                                    run.font.name = fn
                            else:
                                r = p.add_run()
                                r.font.name = fn


def _slides_with_structured_visuals(deck_json: Dict[str, Any]) -> Set[int]:
    """
    deck_json 슬라이드 중 TABLE/DIAGRAM/CHART 스펙이 있는 슬라이드 index(0-based)를 반환한다.
    이런 슬라이드는 Gamma가 그림(PICTURE)으로 렌더링할 수 있어 PICTURE를 보존해야 한다.
    """
    keep: Set[int] = set()
    slides = (deck_json or {}).get("slides") or []
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            continue
        table_md = _norm(s.get("TABLE_MD") or "")
        diagram = _norm(s.get("DIAGRAM_SPEC_KO") or "")
        chart = _norm(s.get("CHART_SPEC_KO") or "")
        if table_md or diagram or chart:
            keep.add(i)
    return keep


def _slide_has_structured_spec(slide_spec: Dict[str, Any]) -> bool:
    if not isinstance(slide_spec, dict):
        return False
    return bool(
        _norm(slide_spec.get("TABLE_MD") or "")
        or _norm(slide_spec.get("DIAGRAM_SPEC_KO") or "")
        or _norm(slide_spec.get("CHART_SPEC_KO") or "")
    )


def _slides_need_generated_image(deck_json: Dict[str, Any]) -> Set[int]:
    need: Set[int] = set()
    slides = (deck_json or {}).get("slides") or []
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            continue
        if bool(s.get("image_needed")):
            need.add(i)
    return need


def _remove_visual_placeholders(
    prs: Presentation,
    keep_picture_slide_idxs: Set[int],
    *,
    keep_placeholder_slide_idxs: Optional[Set[int]] = None,
    remove_pictures: bool = True,
) -> int:
    """
    - AI 이미지/placeholder만 제거
    - 구조화된 시각자료(표/차트/다이어그램 스펙이 있는 슬라이드)의 PICTURE는 보존
    """
    removed = 0
    for si, slide in enumerate(prs.slides):
        for sh in list(slide.shapes):
            # 1) 실제 그림
            if remove_pictures and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 표/다이어그램이 그림으로 들어간 슬라이드는 보존
                if si in keep_picture_slide_idxs:
                    continue
                _remove_shape(sh)
                removed += 1
                continue

            # 2) placeholder (picture)
            if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if keep_placeholder_slide_idxs and si in keep_placeholder_slide_idxs:
                    continue
                ph = getattr(sh, "placeholder_format", None)
                ph_type = str(getattr(ph, "type", "")).lower() if ph else ""
                nm = (getattr(sh, "name", "") or "").lower()
                if ("picture" in ph_type) or ("pic" in ph_type) or ("image" in nm) or ("picture" in nm):
                    _remove_shape(sh)
                    removed += 1
                    continue

            # 3) 도형으로 만든 '이미지 자리'만 제거
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if keep_placeholder_slide_idxs and si in keep_placeholder_slide_idxs:
                    continue
                nm = (getattr(sh, "name", "") or "").lower()
                if any(k in nm for k in ["picture placeholder", "image placeholder", "picture", "image"]):
                    _remove_shape(sh)
                    removed += 1
                    continue

    return removed


def _trim_ending_slides(prs: Presentation) -> None:
    """
    마지막에 감사합니다/Q&A가 여러 장 생기거나, 불필요한 엔딩이 붙는 경우 정리한다.
    - '감사합니다'는 1장만 유지
    - '추가 정보/문의/연락처/회사 소개' 슬라이드는 제거
    """
    bad_keywords = ["추가 정보", "문의", "연락처", "회사 소개", "contact", "thank you", "thanks"]
    thanks_idx: List[int] = []

    for i, slide in enumerate(prs.slides):
        txt = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt += " " + (sh.text_frame.text or "")
        t = _norm(txt).lower()
        if "감사합니다" in t:
            thanks_idx.append(i)

    # '감사합니다'가 2장 이상이면 마지막 1장만 유지
    if len(thanks_idx) > 1:
        for idx in reversed(thanks_idx[:-1]):
            _delete_slide(prs, idx)

    # 마지막 쪽 bad keyword 슬라이드 제거 (보수적으로 마지막 5장 범위)
    n = len(prs.slides)
    for idx in reversed(range(max(0, n - 5), n)):
        slide = prs.slides[idx]
        txt = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt += " " + (sh.text_frame.text or "")
        t = _norm(txt)
        if any(k.lower() in t.lower() for k in bad_keywords):
            # 단 '감사합니다'는 유지
            if "감사합니다" in t:
                continue
            _delete_slide(prs, idx)


def _remove_duplicate_text_shapes(prs: Presentation) -> int:
    removed = 0
    for slide in prs.slides:
        seen: Set[str] = set()
        for sh in list(slide.shapes):
            if not getattr(sh, "has_text_frame", False):
                continue
            t = _norm(sh.text_frame.text or "")
            if not t:
                continue
            key = re.sub(r"\s+", " ", t).strip().lower()
            if len(key) < 10:
                continue
            if key in seen:
                _remove_shape(sh)
                removed += 1
                continue
            seen.add(key)
    return removed


def postprocess_diagrams(pptx_path: str, deck_json: Dict[str, Any], state: Optional[Dict[str, Any]] = None) -> str:
    """
    Gamma 결과 PPTX 후처리
    - 표지/목차 강제 시작부 정리 (1~2번 슬라이드)
    - AI 이미지/placeholder 제거 (단 표/다이어그램 그림은 보존)
    - 엔딩 슬라이드 정리
    """
    prs = Presentation(pptx_path)

    # deck_json 기반으로 "그림을 보존해야 하는 슬라이드" 계산
    keep_picture_slide_idxs = _slides_with_structured_visuals(deck_json)
    need_image_slide_idxs = _slides_need_generated_image(deck_json)
    deck_slides = (deck_json or {}).get("slides") or []

    rewrite_cover = bool((state or {}).get("postprocess_rewrite_cover"))
    rewrite_agenda = bool((state or {}).get("postprocess_rewrite_agenda"))
    style_tables = bool((state or {}).get("postprocess_style_tables"))
    trim_ending = bool((state or {}).get("postprocess_trim_ending"))
    force_rewrite_cover = bool((state or {}).get("force_rewrite_cover"))
    apply_template = bool((state or {}).get("postprocess_apply_template"))
    apply_background = bool((state or {}).get("postprocess_apply_background_image"))
    remove_background = bool((state or {}).get("postprocess_remove_background_image"))
    background_profile = _norm((state or {}).get("postprocess_background_profile") or os.environ.get("POSTPROCESS_BACKGROUND_PROFILE") or "").lower()
    background_image_path = (
        (state or {}).get("postprocess_background_image_path")
        or os.environ.get("POSTPROCESS_BACKGROUND_IMAGE_PATH")
        or ""
    )

    # 1) 표지/목차 시작부 정리 (선택)
    if rewrite_cover and len(prs.slides) >= 1:
        cover_spec = deck_slides[0] if len(deck_slides) >= 1 and isinstance(deck_slides[0], dict) else {}
        if force_rewrite_cover or (not _slide_has_structured_spec(cover_spec)):
            title = _extract_best_title(state or {"deck_json": deck_json})
            _write_cover(prs.slides[0], title)

    if rewrite_agenda and len(prs.slides) >= 2:
        agenda_spec = deck_slides[1] if len(deck_slides) >= 2 and isinstance(deck_slides[1], dict) else {}
        force_rewrite_agenda = bool((state or {}).get("force_rewrite_agenda"))
        if force_rewrite_agenda and not _slide_has_structured_spec(agenda_spec):
            _write_agenda(prs.slides[1])

    # 1.5) 표지/마지막 슬라이드 색상 포인트 보강
    if rewrite_cover and len(prs.slides) >= 1:
        _decorate_cover_slide(prs.slides[0])
    if rewrite_cover and len(prs.slides) >= 1:
        _decorate_thanks_slide(prs.slides[-1])

    # 2) AI image 정리 (placeholder와 그림 삽입 전 대상만 우선 제거)
    _remove_visual_placeholders(
        prs,
        keep_picture_slide_idxs,
        keep_placeholder_slide_idxs=need_image_slide_idxs,
        remove_pictures=True,
    )
    prs.save(pptx_path)

    # 2.2) placeholder 정리 후 이미지 삽입
    try:
        image_paths = maybe_insert_generated_diagrams(pptx_path, deck_json, state=state)
        if state is not None and image_paths:
            state["generated_diagram_images"] = image_paths
    except Exception as e:
        print(f"[WARN] diagram image generation skipped: {e}")
    prs = Presentation(pptx_path)

    # 2.3) 잔여 image placeholder 정리 (생성된 그림은 보존)
    _remove_visual_placeholders(
        prs,
        keep_picture_slide_idxs=keep_picture_slide_idxs.union(need_image_slide_idxs),
        keep_placeholder_slide_idxs=set(),
        remove_pictures=False,
    )
    _remove_duplicate_text_shapes(prs)
    if apply_template:
        _decorate_content_slides(prs)

    # 2.5) table style (선택)
    if style_tables:
        _style_tables(prs)
    _apply_font_name(prs, (state or {}).get("font_name") or "")
    # 3) 엔딩 정리 (선택)
    if trim_ending:
        _trim_ending_slides(prs)

    if remove_background:
        removed = _remove_background_images(prs)
        print(f"[INFO] background image removed: {removed} shapes")

    if apply_background:
        _remove_background_images(prs)
        if background_profile in {"basic", "brown"}:
            n_bg = _apply_background_profile_to_all_slides(
                prs,
                deck_json,
                profile=background_profile,
                state=state,
            )
        else:
            n_bg = _apply_background_image_to_all_slides(prs, str(background_image_path))
        if n_bg <= 0:
            print(f"[WARN] background image skipped: not found or empty path ({background_image_path})")
        else:
            print(
                f"[INFO] background image inserted: {n_bg} slides "
                f"(marker name={BG_MARKER_NAME}, alt={BG_MARKER_ALT})"
            )

    prs.save(pptx_path)
    return pptx_path


def postprocess_diagrams_node(state: Dict[str, Any]) -> Dict[str, Any]:
    pptx_path = state.get("final_ppt_path") or state.get("gamma_ppt_path")
    deck_json = state.get("deck_json") or {}
    if not pptx_path:
        return state
    postprocess_diagrams(pptx_path, deck_json, state=state)
    return state
