from __future__ import annotations

import os
import re
from datetime import datetime
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt

DEFAULT_LAYOUT_WHITELIST = [
    "Title Slide",
    "Title and Content",
    "Two Content",
]

LAYOUT_ID_TO_ROLE = {
    "cover": "cover",
    "agenda": "text",
    "text": "text",
    "content": "content",
    "two_content": "content",
    "qna": "text",
}


def _norm(s: Any) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()


def _safe_filename(name: str) -> str:
    name = re.sub(r"[\\/:*?\"<>|]+", " ", str(name or ""))
    name = re.sub(r"\s+", " ", name).strip()
    if not name:
        return "result"
    return name[:48].rstrip()


def _avoid_windows_lock(path: str) -> str:
    base, ext = os.path.splitext(path)
    if not os.path.exists(path):
        return path
    for i in range(1, 200):
        cand = f"{base} ({i}){ext}"
        if not os.path.exists(cand):
            return cand
    return f"{base}_{int(datetime.now().timestamp())}{ext}"


def _pick_layout(prs: Presentation, preferred_idx: int) -> Any:
    if len(prs.slide_layouts) == 0:
        raise RuntimeError("No slide layouts in template presentation.")
    idx = min(max(preferred_idx, 0), len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


def _layout_name_map(prs: Presentation) -> Dict[str, int]:
    out: Dict[str, int] = {}
    for i, ly in enumerate(prs.slide_layouts):
        out[_norm(getattr(ly, "name", "")).lower()] = i
    return out


def _find_layout_idx(prs: Presentation, keywords: List[str], fallback: int) -> int:
    name_map = _layout_name_map(prs)
    for key in keywords:
        key_l = key.lower()
        for n, i in name_map.items():
            if key_l in n:
                return i
    return min(max(fallback, 0), len(prs.slide_layouts) - 1)


def _resolve_allowed_layout_indices(prs: Presentation, whitelist_names: List[str]) -> List[int]:
    if not whitelist_names:
        return list(range(len(prs.slide_layouts)))
    allowed: List[int] = []
    norm_w = [_norm(x).lower() for x in whitelist_names if _norm(x)]
    for i, ly in enumerate(prs.slide_layouts):
        nm = _norm(getattr(ly, "name", "")).lower()
        if nm in norm_w:
            allowed.append(i)
    return allowed


def _pick_first_allowed(allowed_indices: List[int], preferred: int) -> int:
    if preferred in allowed_indices:
        return preferred
    if allowed_indices:
        return allowed_indices[0]
    return preferred


def _find_placeholder(slide, types: List[int], *, prefer_idx: Optional[int] = None):
    cands = []
    for ph in slide.placeholders:
        pf = getattr(ph, "placeholder_format", None)
        if pf is None:
            continue
        pht = int(getattr(pf, "type", -1))
        idx = int(getattr(pf, "idx", -1))
        if pht in types and getattr(ph, "has_text_frame", False):
            area = int(getattr(ph, "width", 0)) * int(getattr(ph, "height", 0))
            cands.append((ph, idx, area))
    if not cands:
        return None
    if prefer_idx is not None:
        for ph, idx, _ in cands:
            if idx == prefer_idx:
                return ph
    cands.sort(key=lambda x: x[2], reverse=True)
    return cands[0][0]


def _bump(stats: Dict[str, int], key: str, inc: int = 1) -> None:
    stats[key] = int(stats.get(key, 0)) + inc


def _set_title(slide, title: str, *, strict_placeholder_only: bool, stats: Dict[str, int]) -> None:
    title_box = slide.shapes.title or _find_placeholder(
        slide,
        [int(PP_PLACEHOLDER.TITLE), int(PP_PLACEHOLDER.CENTER_TITLE)],
    )
    if title_box is not None:
        title_box.text = title
        _bump(stats, "title_placeholder")
        return
    _bump(stats, "title_missing")
    if strict_placeholder_only:
        return
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.8))
    t.text_frame.text = title
    _bump(stats, "textbox_fallback")


def _fill_text_placeholder(ph, lines: List[str]) -> None:
    tf = ph.text_frame
    tf.clear()
    if not lines:
        return
    p0 = tf.paragraphs[0]
    p0.text = lines[0]
    p0.level = 0
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0


def _add_title_and_body(
    slide,
    title: str,
    body_lines: List[str],
    *,
    strict_placeholder_only: bool,
    stats: Dict[str, int],
) -> None:
    _set_title(slide, title, strict_placeholder_only=strict_placeholder_only, stats=stats)
    body_placeholder = _find_placeholder(
        slide,
        [
            int(PP_PLACEHOLDER.BODY),
            int(PP_PLACEHOLDER.OBJECT),
            int(PP_PLACEHOLDER.VERTICAL_BODY),
            int(PP_PLACEHOLDER.VERTICAL_OBJECT),
        ],
        prefer_idx=1,
    )
    if body_placeholder is not None and getattr(body_placeholder, "has_text_frame", False):
        _fill_text_placeholder(body_placeholder, body_lines)
        _bump(stats, "body_placeholder")
        return
    _bump(stats, "body_missing")
    if strict_placeholder_only:
        return
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.4))
    _fill_text_placeholder(body, body_lines)
    _bump(stats, "textbox_fallback")


def _add_title_two_content(
    slide,
    title: str,
    left_lines: List[str],
    right_lines: List[str],
    *,
    strict_placeholder_only: bool,
    stats: Dict[str, int],
) -> None:
    _set_title(slide, title, strict_placeholder_only=strict_placeholder_only, stats=stats)
    left_ph = _find_placeholder(slide, [int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)], prefer_idx=1)
    right_ph = _find_placeholder(slide, [int(PP_PLACEHOLDER.BODY), int(PP_PLACEHOLDER.OBJECT)], prefer_idx=2)
    if left_ph is not None:
        _fill_text_placeholder(left_ph, left_lines)
        _bump(stats, "body_placeholder")
    if right_ph is not None:
        _fill_text_placeholder(right_ph, right_lines)
        _bump(stats, "body_placeholder")
    if left_ph is not None or right_ph is not None:
        return
    _bump(stats, "body_missing")
    if strict_placeholder_only:
        return
    _add_title_and_body(
        slide,
        title,
        left_lines + right_lines,
        strict_placeholder_only=strict_placeholder_only,
        stats=stats,
    )


def _is_cover_section(sec: str) -> bool:
    return sec in {"표지", "cover"}


def _is_agenda_section(sec: str) -> bool:
    return sec in {"목차", "agenda"}


def _is_content_heavy_section(sec: str) -> bool:
    return sec in {"연구 내용", "추진 계획", "사업화 전략 및 계획", "content"}


def _pick_layout_index_by_section(prs: Presentation, section: str, use_two_col: bool) -> int:
    sec = _norm(section)
    if _is_cover_section(sec):
        return _find_layout_idx(prs, ["title slide"], 0)
    if _is_agenda_section(sec):
        return _find_layout_idx(prs, ["title and content", "section header"], 1 if len(prs.slide_layouts) > 1 else 0)
    if sec in {"Q&A", "q&a", "qna"}:
        return _find_layout_idx(prs, ["section header", "title only"], 2 if len(prs.slide_layouts) > 2 else 0)
    if use_two_col or _is_content_heavy_section(sec):
        return _find_layout_idx(prs, ["two content", "comparison"], 2 if len(prs.slide_layouts) > 2 else 0)
    return _find_layout_idx(prs, ["title and content", "content with caption"], 1 if len(prs.slide_layouts) > 1 else 0)


def _slide_body_lines(slide_spec: Dict[str, Any]) -> List[str]:
    out: List[str] = []
    key = _norm(slide_spec.get("key_message"))
    if key:
        out.append(key)
    for b in (slide_spec.get("bullets") or [])[:6]:
        bt = _norm(b)
        if bt:
            out.append(f"- {bt}")
    if not out:
        out.append("- Content summary")
    return out


def _pick_layout_index_for_slide(
    prs: Presentation,
    slide_spec: Dict[str, Any],
    *,
    allowed_indices: List[int],
    cover_mode: bool,
) -> int:
    section = _norm(slide_spec.get("section"))
    use_two_col = _is_content_heavy_section(section) or len(_slide_body_lines(slide_spec)) >= 6
    if cover_mode:
        fallback = _pick_layout_index_by_section(prs, "표지", False)
    else:
        fallback = _pick_layout_index_by_section(prs, section, use_two_col)

    layout_id = _norm(slide_spec.get("layout_id") or slide_spec.get("slide_layout")).lower()
    role = LAYOUT_ID_TO_ROLE.get(layout_id, "")
    if role == "cover":
        return _pick_first_allowed(allowed_indices, _find_layout_idx(prs, ["title slide"], fallback))
    if role == "content":
        return _pick_first_allowed(allowed_indices, _find_layout_idx(prs, ["two content", "comparison"], fallback))
    if role == "text":
        return _pick_first_allowed(allowed_indices, _find_layout_idx(prs, ["title and content", "section header"], fallback))
    return _pick_first_allowed(allowed_indices, fallback)


def _parse_table_md(table_md: str) -> List[List[str]]:
    rows: List[List[str]] = []
    for raw in (table_md or "").splitlines():
        line = raw.strip()
        if not line or "|" not in line:
            continue
        if re.fullmatch(r"\|?[\-\s:|]+\|?", line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        if any(cells):
            rows.append(cells)
    return rows


def _table_rows_to_lines(rows: List[List[str]]) -> List[str]:
    if not rows:
        return []
    out: List[str] = []
    header = [c for c in rows[0] if c]
    if header:
        out.append(" | ".join(header))
    for r in rows[1:7]:
        vals = [c for c in r if c]
        if vals:
            out.append("- " + " / ".join(vals[:3]))
    return out


def _add_table(slide, rows: List[List[str]]) -> bool:
    if not rows:
        return False
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    left, top, width, height = Inches(0.9), Inches(1.8), Inches(11.4), Inches(4.8)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    for r in range(n_rows):
        for c in range(n_cols):
            txt = rows[r][c] if c < len(rows[r]) else ""
            cell = table.cell(r, c)
            cell.text = txt
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12 if r == 0 else 11)
                    if r == 0:
                        run.font.bold = True
    return True


def template_render_node(state: Dict[str, Any]) -> Dict[str, Any]:
    deck = state.get("deck_json") or {}
    slides = deck.get("slides") or []
    if not slides:
        raise RuntimeError("deck_json.slides is empty. Nothing to render.")

    strict_placeholder_only = bool(state.get("template_strict_placeholder_only", True))
    table_as_shape = bool(state.get("template_table_as_shape", False))
    whitelist_names = state.get("template_layout_whitelist") or DEFAULT_LAYOUT_WHITELIST

    template_path = _norm(state.get("template_pptx_path") or state.get("template_ppt_path"))
    if not template_path or not os.path.exists(template_path):
        raise RuntimeError(f"Template PPTX not found: {template_path or '(empty)'}")
    prs = Presentation(template_path)

    print(f"[TEMPLATE] template_path={template_path}")
    for i, ly in enumerate(prs.slide_layouts):
        print(f"[TEMPLATE] layout[{i}]={_norm(getattr(ly, 'name', ''))}")

    allowed_indices = _resolve_allowed_layout_indices(prs, whitelist_names)
    if not allowed_indices:
        raise RuntimeError(
            f"template_layout_whitelist not matched: {whitelist_names}. "
            "Please set exact layout names in state.template_layout_whitelist."
        )
    print(f"[TEMPLATE] whitelist={whitelist_names}")
    print(f"[TEMPLATE] allowed_layout_indices={allowed_indices}")

    while len(prs.slides) > 0:
        sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
        sld = sld_id_lst[0]
        rel_id = sld.rId
        sld_id_lst.remove(sld)
        prs.part.drop_rel(rel_id)

    stats: Dict[str, int] = {}

    for i, s in enumerate(slides):
        sec = _norm(s.get("section"))
        title = _norm(s.get("slide_title")) or "Slide"

        if i == 0 or _is_cover_section(sec):
            layout_idx = _pick_layout_index_for_slide(prs, s, allowed_indices=allowed_indices, cover_mode=True)
            slide = prs.slides.add_slide(_pick_layout(prs, layout_idx))
            _set_title(slide, title, strict_placeholder_only=strict_placeholder_only, stats=stats)
            sub = _find_placeholder(slide, [int(PP_PLACEHOLDER.SUBTITLE)], prefer_idx=1)
            cover_sub = _norm(deck.get("deck_title")) or "R&D proposal deck"
            if sub is not None:
                _fill_text_placeholder(sub, [cover_sub])
                _bump(stats, "body_placeholder")
            else:
                _add_title_and_body(
                    slide,
                    title,
                    [cover_sub],
                    strict_placeholder_only=strict_placeholder_only,
                    stats=stats,
                )
            continue

        if _is_agenda_section(sec):
            layout_idx = _pick_layout_index_for_slide(prs, s, allowed_indices=allowed_indices, cover_mode=False)
            slide = prs.slides.add_slide(_pick_layout(prs, layout_idx))
            body = [f"- {x}" for x in (_norm(s.get("TABLE_MD")) or "").splitlines() if x.strip() and "|" not in x]
            if not body:
                body = [f"- {x}" for x in (s.get("bullets") or [])]
            _add_title_and_body(
                slide,
                "목차",
                body[:10],
                strict_placeholder_only=strict_placeholder_only,
                stats=stats,
            )
            continue

        use_two_col = _is_content_heavy_section(sec) or len(_slide_body_lines(s)) >= 6
        layout_idx = _pick_layout_index_for_slide(prs, s, allowed_indices=allowed_indices, cover_mode=False)
        slide = prs.slides.add_slide(_pick_layout(prs, layout_idx))
        rows = _parse_table_md(str(s.get("TABLE_MD") or ""))
        lines = _slide_body_lines(s)

        if rows and not use_two_col:
            _set_title(slide, title, strict_placeholder_only=strict_placeholder_only, stats=stats)
            if strict_placeholder_only or not table_as_shape:
                _add_title_and_body(
                    slide,
                    title,
                    _table_rows_to_lines(rows) or lines,
                    strict_placeholder_only=strict_placeholder_only,
                    stats=stats,
                )
            else:
                if not _add_table(slide, rows):
                    _add_title_and_body(
                        slide,
                        title,
                        lines,
                        strict_placeholder_only=strict_placeholder_only,
                        stats=stats,
                    )
        elif use_two_col:
            left = lines[::2]
            right = lines[1::2] or ["- Details"]
            if rows:
                right = [f"- {', '.join(r[:2]).strip()}" for r in rows[1:6] if any(r)] or right
            _add_title_two_content(
                slide,
                title,
                left[:6],
                right[:6],
                strict_placeholder_only=strict_placeholder_only,
                stats=stats,
            )
        else:
            _add_title_and_body(
                slide,
                title,
                lines,
                strict_placeholder_only=strict_placeholder_only,
                stats=stats,
            )

    output_dir = _norm(state.get("output_dir") or "output")
    os.makedirs(output_dir, exist_ok=True)
    if _norm(state.get("output_filename")):
        out_name = _norm(state.get("output_filename"))
    else:
        ts = datetime.now().strftime("%Y%m%d")
        base = _safe_filename(_norm(deck.get("deck_title")) or "result")
        out_name = f"RanDi_{base}_{ts}_template.pptx"
    out_path = _avoid_windows_lock(os.path.join(output_dir, out_name))
    prs.save(out_path)

    print(
        "[TEMPLATE] placeholder stats:",
        f"title={stats.get('title_placeholder', 0)}",
        f"body={stats.get('body_placeholder', 0)}",
        f"title_missing={stats.get('title_missing', 0)}",
        f"body_missing={stats.get('body_missing', 0)}",
        f"textbox_fallback={stats.get('textbox_fallback', 0)}",
    )

    state["template_ppt_path"] = out_path
    state["final_ppt_path"] = out_path
    return state
